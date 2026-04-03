import os
import re
import tempfile
import fitz  # PyMuPDF
from supabase import create_client
from .config import settings
from .vector_store import add_documents

STORAGE_BUCKET = "thecrowsnest"

_supabase = create_client(settings.SUPABASE_URL, settings.SUPABASE_SERVICE_ROLE_KEY)

# ============================================================
# Extraction-gate thresholds (mirrors upload-safety.ts)
# ============================================================

MIN_EXTRACT_CHARS = 120
MIN_OCR_CHARS = 80
MAX_GARBLE_RATIO = 0.35
MAX_PAGES_AUTO = 300


def _garble_ratio(text: str) -> float:
    """Fraction of non-printable / non-standard characters in extracted text."""
    if not text:
        return 1.0
    garbled = sum(1 for ch in text if not ch.isprintable() and ch not in "\n\r\t")
    return garbled / len(text)


def download_and_extract_text(storage_key: str) -> list[str]:
    fd, local_path = tempfile.mkstemp(suffix=os.path.splitext(storage_key)[1] or ".pdf")
    os.close(fd)

    print(f"Downloading {storage_key} from Supabase Storage to {local_path}...")
    file_bytes = _supabase.storage.from_(STORAGE_BUCKET).download(storage_key)
    with open(local_path, "wb") as f:
        f.write(file_bytes)

    # Extract text
    texts = []
    if local_path.lower().endswith(".pdf"):
        try:
            doc = fitz.open(local_path)
            for page in doc:
                texts.append(page.get_text())
            doc.close()
        except Exception as e:
            print(f"PyMuPDF fast-extraction failed: {e}. Activating pdfplumber Dual-Engine fallback...")
            import pdfplumber
            with pdfplumber.open(local_path) as pdf:
                for page in pdf.pages:
                    extracted = page.extract_text()
                    texts.append(extracted or "")
    elif local_path.lower().endswith(".pptx"):
        from pptx import Presentation
        prs = Presentation(local_path)
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            texts.append("\n".join(slide_text))
    elif local_path.lower().endswith(".docx"):
        from docx import Document
        doc = Document(local_path)
        content = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        texts = [content[i:i+2000] for i in range(0, len(content), 2000)]
    elif local_path.lower().endswith((".png", ".jpg", ".jpeg", ".webp")):
        from .ai import extract_text_from_image
        content = extract_text_from_image(local_path)
        if content:
            texts = [content[i:i+2000] for i in range(0, len(content), 2000)]
    else:
        with open(local_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
            texts = [content[i:i+2000] for i in range(0, len(content), 2000)]

    # Cleanup
    if os.path.exists(local_path):
        os.remove(local_path)

    # Strip null bytes to prevent PostgreSQL 'cannot contain NUL' errors
    return [t.replace('\x00', '') for t in texts if t.strip()]


def _compute_extraction_metrics(texts: list[str], storage_key: str) -> dict:
    """Compute quality metrics from extracted text pages."""
    full_text = " ".join(texts)
    is_image = storage_key.lower().endswith((".png", ".jpg", ".jpeg", ".webp"))

    return {
        "extractCharCount": len(full_text),
        "pageCount": len(texts),
        "ocrUsed": is_image,
        "garbleRatio": round(_garble_ratio(full_text), 4),
        "chunkCount": 0,  # filled later by the embedding step
    }


def _run_extraction_gate(texts: list[str], metrics: dict) -> dict | None:
    """
    Deterministic extraction-gate checks (EX-001 through EX-006).
    Returns a rejection dict if a gate fails, or None if all pass.
    """
    # EX-001: No text at all (parser couldn't open)
    if not texts:
        return {
            "evaluation": "REJECTED",
            "confidence": 0,
            "reason": "File could not be read.",
            "reasonCode": "unreadable_file",
            "metrics": metrics,
        }

    # EX-003: Excessive pages → route to review, don't reject
    if metrics["pageCount"] > MAX_PAGES_AUTO:
        return {
            "evaluation": "PENDING",
            "confidence": 0,
            "reason": f"Document has {metrics['pageCount']} pages and needs admin review.",
            "reasonCode": "excessive_pages",
            "metrics": metrics,
        }

    # EX-004: Too little text
    if metrics["extractCharCount"] < MIN_EXTRACT_CHARS:
        return {
            "evaluation": "REJECTED",
            "confidence": 0,
            "reason": "Not enough readable text found.",
            "reasonCode": "empty_or_low_text",
            "metrics": metrics,
        }

    # EX-005: Garbled text
    if metrics["garbleRatio"] > MAX_GARBLE_RATIO:
        return {
            "evaluation": "REJECTED",
            "confidence": 0,
            "reason": "Extracted text appears corrupted.",
            "reasonCode": "garbled_text",
            "metrics": metrics,
        }

    # EX-006: Image OCR too little text
    if metrics["ocrUsed"] and metrics["extractCharCount"] < MIN_OCR_CHARS:
        return {
            "evaluation": "REJECTED",
            "confidence": 0,
            "reason": "Not enough readable text found in image.",
            "reasonCode": "low_ocr_content",
            "metrics": metrics,
        }

    return None  # all gates passed


def process_material(class_id: str, material_id: str, storage_key: str, file_name: str):
    print(f"Processing material: {file_name} for class: {class_id}")
    texts = download_and_extract_text(storage_key)
    metadatas = [{"source": file_name, "page": i+1} for i in range(len(texts))]
    add_documents(class_id, material_id, texts, metadatas)
    print("Added material to PostgreSQL Database.")


def download_extract_and_evaluate(class_id: str, material_id: str, storage_key: str, file_name: str, class_context: str) -> dict:
    """Autonomous evaluation logic with extraction gate."""
    print(f"Autonomous Evaluation for material: {file_name} (class: {class_id})")

    # 1. Extract text
    try:
        texts = download_and_extract_text(storage_key)
    except Exception as e:
        print(f"Extraction failed: {e}")
        return {
            "evaluation": "REJECTED",
            "confidence": 0,
            "reason": "File could not be read.",
            "reasonCode": "unreadable_file",
            "metrics": {"extractCharCount": 0, "pageCount": 0, "ocrUsed": False, "garbleRatio": 1.0, "chunkCount": 0},
        }

    # 2. Compute extraction metrics
    metrics = _compute_extraction_metrics(texts, storage_key)

    # 3. Run extraction gate (deterministic checks before AI call)
    gate_result = _run_extraction_gate(texts, metrics)
    if gate_result is not None:
        print(f"Extraction gate failed: {gate_result['reasonCode']}")
        return gate_result

    # 4. Grab a snippet of the text (first ~15,000 characters to fit context window comfortably)
    full_text = " ".join(texts)
    snippet = full_text[:15000]

    # 5. Evaluate using Google Gemini
    from .ai import evaluate_material_against_syllabus
    try:
        results = evaluate_material_against_syllabus(class_context, snippet)
        confidence = results.get("confidence", 50)
        reason = results.get("reason", "Unknown reason.")
        reason_code = results.get("reasonCode")  # may be provided by deterministic checks in ai.py
    except Exception as e:
        print(f"AI Evaluation failed: {e}")
        return {
            "evaluation": "PENDING",
            "confidence": 50,
            "reason": "AI evaluation crashed. Needs manual review.",
            "reasonCode": "uncertain_relevance",
            "metrics": metrics,
        }

    print(f"Evaluation resulted in confidence {confidence}: {reason}")

    # 6. Tri-State Logic Execution
    if confidence >= 75:
        print(f"Confidence {confidence} >= 75. Approving.")
        metadatas = [{"source": file_name, "page": i+1} for i in range(len(texts))]
        metrics["chunkCount"] = len(texts)
        return {
            "evaluation": "APPROVED",
            "confidence": confidence,
            "reason": reason,
            "reasonCode": reason_code or "approved",
            "metrics": metrics,
            "texts": texts,
            "metadatas": metadatas,
        }
    elif confidence >= 50:
        print(f"Confidence {confidence} between 50 and 74. Leaving as PENDING_REVIEW.")
        return {
            "evaluation": "PENDING",
            "confidence": confidence,
            "reason": reason,
            "reasonCode": reason_code or "uncertain_relevance",
            "metrics": metrics,
        }
    else:
        print("Confidence < 50. Marking REJECTED for auto-deletion.")
        return {
            "evaluation": "REJECTED",
            "confidence": confidence,
            "reason": reason,
            "reasonCode": reason_code or "low_relevance_confidence",
            "metrics": metrics,
        }
